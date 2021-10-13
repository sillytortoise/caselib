import sys

from pptx import Presentation

prs = Presentation("竞品分析模板.pptx")
slide = prs.slides[0]
for shape in slide.placeholders:
    phf = shape.placeholder_format
    print(phf.idx, '\t', shape.text_frame.text)
