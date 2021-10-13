import sys

from pptx import Presentation
import json


if sys.argv[3] == 2:  # 特色化案例库
    prs = Presentation(r"特色化案例库模板.pptx")
    with open('user/%s/%s.json' % (sys.argv[1], sys.argv[2]), 'r', encoding="utf-8") as myjson:
        pages = json.load(myjson)
        for page in pages:
            m = {}
            if page["pics"] is not None:
                for p in page["pics"]:
                    m[p["order"]+21] = p["title"]
                    m[p["order"]+26] = p["path"]

            slide = prs.slides.add_slide(prs.slide_layouts[0])
            for shape in slide.placeholders:
                phf = shape.placeholder_format
                if phf.idx == 10:
                    shape.text_frame.text = page["case_num"]
                elif phf.idx == 11:
                    shape.text_frame.text = str(page["product_class"])
                elif phf.idx == 12:
                    shape.text_frame.text = str(page["feature"])
                elif phf.idx == 13:
                    shape.text_frame.text = str(page["page_num"])
                elif phf.idx == 14:
                    shape.text_frame.text = str(page["name"])
                elif phf.idx == 15:
                    shape.text_frame.text = str(page["version"])
                elif phf.idx == 16:
                    shape.text_frame.text = str(page["app_class"])
                elif phf.idx == 17:
                    shape.text_frame.text = str(page["date"])
                elif phf.idx == 18:
                    shape.text_frame.text = str(page["username"])
                elif phf.idx == 19:
                    shape.text_frame.text = str(page["abstract"])
                elif phf.idx == 20:
                    shape.text_frame.text = str(page["idea"])
                elif phf.idx == 32:
                    shape.text_frame.text = str(page["detail"])
                if page["pics"] is None:
                    continue
                if 22 <= phf.idx <= 26:
                    if phf.idx in m.keys():
                        shape.text_frame.text = m[phf.idx]
                    else:
                        sp = shape.element
                        sp.getparent().remove(sp)
                elif 27 <= phf.idx <= 31:
                    if phf.idx in m.keys():
                        shape.insert_picture(
                            'static/upload/'+m[phf.idx])
                    else:
                        sp = shape.element
                        sp.getparent().remove(sp)

        prs.save('user/%s/%s.pptx' % (sys.argv[1], sys.argv[2]))
else:  # 竞品分析
    prs = Presentation(r"竞品分析模板.pptx")
    with open('user/%s/%s.json' % (sys.argv[1], sys.argv[2]), 'r', encoding="utf-8") as myjson:
        pages = json.load(myjson)
        for page in pages:
            m = {}
            if page["pics"] is not None:
                for p in page["pics"]:
                    m[p["order"]+19] = p["title"]
                    m[p["order"]+24] = p["path"]

            print(m)
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            for shape in slide.placeholders:
                phf = shape.placeholder_format

                if phf.idx == 10:
                    shape.text_frame.text = str(
                        page["target"])+'/'+str(page["item"])
                elif phf.idx == 12:
                    shape.text_frame.text = str(page["target_content"])
                elif phf.idx == 14:
                    shape.text_frame.text = str(page["priority"])
                elif phf.idx == 30:
                    shape.text_frame.text = str(page["type"])
                elif phf.idx == 17:
                    shape.text_frame.text = str(page["problem"])
                elif phf.idx == 19:
                    shape.text_frame.text = str(page["advice"])
                if page["pics"] is None:
                    continue
                if 20 <= phf.idx <= 24:
                    if phf.idx in m.keys():
                        shape.text_frame.text = '同业参考：'+m[phf.idx]
                    else:
                        sp = shape.element
                        sp.getparent().remove(sp)
                elif 25 <= phf.idx <= 29:
                    if phf.idx in m.keys():
                        shape.insert_picture(
                            'static/upload/'+m[phf.idx])
                    else:
                        sp = shape.element
                        sp.getparent().remove(sp)

        prs.save('user/%s/%s.pptx' % (sys.argv[1], sys.argv[2]))
